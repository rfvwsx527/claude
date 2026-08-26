# -*- coding: utf-8 -*-
"""在既有簡報後面接兩張案例投影片。
   所有形狀都從原本的投影片複製而來，字體／顏色／圓角外框因此與前三張一致。"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

AMBER, WHITE, BODY, FAINT = "D9A441", "EEF3FA", "93A4BC", "6B7B95"
GREEN, BLUE, PURPLE, RED, ORANGE = "35C79A", "6E8CFF", "A78BFA", "F0576B", "E8A33D"
CJK, MONO = "Microsoft JhengHei", "Consolas"

def by_name(slide):
    d = {}
    for sh in slide.shapes:
        d.setdefault(sh.name, []).append(sh)
    return d

def drop(sh):
    sh._element.getparent().remove(sh._element)

def place(sh, l, t, w, h=None):
    sh.left, sh.top, sh.width = Inches(l), Inches(t), Inches(w)
    if h is not None:
        sh.height = Inches(h)

def _max_id(slide):
    return max((int(e.get("id")) for e in slide.shapes._spTree.iter(qn("p:cNvPr"))), default=1)

def clone(slide, sh, l, t, w, h=None, name=None):
    el = copy.deepcopy(sh._element)
    slide.shapes._spTree.append(el)
    new = slide.shapes[-1]
    cNvPr = el.find(".//" + qn("p:cNvPr"))
    cNvPr.set("id", str(_max_id(slide) + 1))
    if name:
        cNvPr.set("name", name)
    place(new, l, t, w, h)
    return new

def style(run, size, bold, color, font=CJK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = parse_xml('<a:%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>' % tag.split(":")[1])
            rPr.append(el)
        el.set("typeface", font)

def write(sh, lines, space_after=0, align=None):
    """lines: [(text, size, bold, color, font, space_after_pt)] — 沿用第一段的 pPr"""
    tf = sh.text_frame
    tf.word_wrap = True
    body = tf._txBody
    p_tpl = copy.deepcopy(tf.paragraphs[0]._p)
    for r in p_tpl.findall(qn("a:r")):
        p_tpl.remove(r)
    for br in p_tpl.findall(qn("a:br")):
        p_tpl.remove(br)
    for p in body.findall(qn("a:p")):
        body.remove(p)
    from pptx.text.text import _Paragraph
    for spec in lines:
        text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
        font = spec[4] if len(spec) > 4 else CJK
        sa = spec[5] if len(spec) > 5 else space_after
        p_el = copy.deepcopy(p_tpl)
        body.append(p_el)
        para = _Paragraph(p_el, tf)
        run = para.add_run()
        run.text = text
        style(run, size, bold, color, font)
        if sa:
            para.space_after = Pt(sa)
        if align is not None:
            para.alignment = align
    return sh

def mixed(sh, runs):
    """同一段落內放多個 run，用於中英混排的標籤"""
    tf = sh.text_frame
    tf.word_wrap = True
    body = tf._txBody
    p_tpl = copy.deepcopy(tf.paragraphs[0]._p)
    for r in p_tpl.findall(qn("a:r")):
        p_tpl.remove(r)
    for p in body.findall(qn("a:p")):
        body.remove(p)
    from pptx.text.text import _Paragraph
    p_el = copy.deepcopy(p_tpl)
    body.append(p_el)
    para = _Paragraph(p_el, tf)
    for text, size, bold, color, font in runs:
        r = para.add_run()
        r.text = text
        style(r, size, bold, color, font)
    return sh

prs = Presentation("work.pptx")
s4, s5 = prs.slides[3], prs.slides[4]

# ─────────────────────────── 第四張：案例 ───────────────────────────
n = by_name(s4)
for name in ("Shape 13", "Image 1", "Text 11", "Text 12", "Text 14", "Text 15"):
    for sh in n.get(name, []):
        drop(sh)
big_panel, img_frame = n["Shape 10"][0], n["Shape 10"][1]
drop(big_panel)                             # 右側大面板，這張不需要

mixed(n["Text 0"][0], [("CASE STUDY", 10, False, AMBER, MONO), ("  ·  實際做出來的東西", 10, False, AMBER, CJK)])
place(n["Text 0"][0], 0.62, 0.46, 9.0)

t1 = n["Text 1"][0]
place(t1, 0.62, 0.74, 10.6, 0.62)
tf = t1.text_frame
body = tf._txBody
p_tpl = copy.deepcopy(tf.paragraphs[0]._p)
for r in p_tpl.findall(qn("a:r")):
    p_tpl.remove(r)
for p in body.findall(qn("a:p")):
    body.remove(p)
from pptx.text.text import _Paragraph
p_el = copy.deepcopy(p_tpl); body.append(p_el)
para = _Paragraph(p_el, tf)
for text, color in (("技能一多就找不到，", WHITE), ("那就先做一個索引", AMBER)):
    r = para.add_run(); r.text = text; style(r, 30, True, color)

# 左欄：痛點
place(n["Shape 2"][0], 0.62, 1.60, 3.60, 3.42)
write(n["Text 3"][0], [("01", 14, True, AMBER)]);            place(n["Text 3"][0], 0.86, 1.77, 0.40, 0.28)
write(n["Text 4"][0], [("痛點", 15, True, WHITE)]);           place(n["Text 4"][0], 1.26, 1.78, 3.08, 0.32)
write(n["Text 5"][0], [("PAIN POINTS", 8, False, FAINT, MONO)]); place(n["Text 5"][0], 0.88, 2.38, 3.08, 0.20)
write(n["Text 6"][0], [
    ("沒有專屬的分類介面", 12, True, WHITE, CJK, 3),
    ("Claude 目前沒有地方可以把技能分門別類、快速查找，只能靠記憶。", 11, False, BODY, CJK, 14),
    ("技能一多就得逐一點開", 12, True, WHITE, CJK, 3),
    ("四十個技能散在範例、內建、自建三個目錄，想知道某個技能做什麼，"
     "只能一個一個開 SKILL.md 看。", 11, False, BODY, CJK, 0),
])
place(n["Text 6"][0], 0.88, 2.62, 3.10, 2.28)   # 收在面板內，別壓到下方數字

# 右欄：成果
write(n["Text 7"][0], [("02", 14, True, GREEN)]);   place(n["Text 7"][0], 4.52, 1.75, 0.40, 0.28)
write(n["Text 8"][0], [("成果", 15, True, WHITE)]);  place(n["Text 8"][0], 4.92, 1.77, 1.00, 0.28)
write(n["Text 9"][0], [("用 Claude Artifacts 做的技能庫儀表板 — 分類、搜尋、篩選，開啟就能用",
                        10, False, BODY)])
place(n["Text 9"][0], 6.05, 1.80, 6.66, 0.24)

place(img_frame, 4.52, 2.20, 8.19, 4.15)    # 沿用原本的圖片外框樣式
img = n["Image 0"][0]
img._element.getparent().remove(img._element)
s4.shapes.add_picture("../dashboard.png", Inches(4.58), Inches(2.26), Inches(8.07), Inches(4.03))

# 現況數字：填滿左欄下方，同時提供具體規模
figs_panel = clone(s4, n["Shape 2"][0], 0.62, 5.22, 3.60, 1.13, "Shape figures")
for i, (v, k) in enumerate((("40", "技能"), ("10", "分類"), ("3", "個來源"))):
    x = 0.86 + i * 1.06
    fv = clone(s4, n["Text 4"][0], x, 5.40, 1.00, 0.42, "Text fig%d v" % i)  # 22pt 需要 0.42"
    write(fv, [(v, 22, True, WHITE)])
    fk = clone(s4, n["Text 5"][0], x, 5.88, 1.00, 0.20, "Text fig%d k" % i)
    write(fk, [(k, 8.5, False, FAINT)])

write(n["Text 16"][0], [("案例 · 技能庫儀表板", 9, False, FAINT)])

# ─────────────────────────── 第五張：製作流程 ───────────────────────────
m = by_name(s5)
for name in ("Shape 13", "Image 1", "Text 11", "Text 12", "Text 14", "Text 15",
             "Text 3", "Text 4", "Text 5", "Text 6", "Text 9"):
    for sh in m.get(name, []):
        drop(sh)
for sh in m["Shape 10"]:                    # 大面板與圖片框都不要
    drop(sh)
drop(m["Image 0"][0])

mixed(m["Text 0"][0], [("HOW IT WAS BUILT", 10, False, AMBER, MONO), ("  ·  兩個步驟", 10, False, AMBER, CJK)])
place(m["Text 0"][0], 0.62, 0.46, 9.0)

t1b = m["Text 1"][0]
place(t1b, 0.62, 0.74, 10.6, 0.62)
tf = t1b.text_frame; body = tf._txBody
p_tpl = copy.deepcopy(tf.paragraphs[0]._p)
for r in p_tpl.findall(qn("a:r")):
    p_tpl.remove(r)
for p in body.findall(qn("a:p")):
    body.remove(p)
p_el = copy.deepcopy(p_tpl); body.append(p_el)
para = _Paragraph(p_el, tf)
for text, color in (("說一次，", WHITE), ("之後每天自己跑", AMBER)):
    r = para.add_run(); r.text = text; style(r, 30, True, color)

# 兩個步驟結構一致：標題在框外，內容在框內
card = m["Shape 2"][0]
place(card, 0.62, 2.20, 12.09, 0.72)        # 步驟一的內容框：只裝 prompt

lbl1 = clone(s5, m["Text 7"][0], 0.90, 1.55, 0.40, 0.28, "Text S1 num")
write(lbl1, [("01", 14, True, AMBER)])
ttl1 = clone(s5, m["Text 8"][0], 1.30, 1.56, 6.00, 0.28, "Text S1 title")
write(ttl1, [("跟 AI 說要做什麼", 15, True, WHITE)])
lab1 = clone(s5, m["Text 0"][0], 1.30, 1.92, 4.00, 0.20, "Text S1 label")
write(lab1, [("PROMPT", 8, False, FAINT, MONO)])
bd1 = clone(s5, m["Text 8"][0], 0.92, 2.42, 11.45, 0.52, "Text S1 body")
write(bd1, [("「用 Claude Artifacts 做一個技能庫儀表板，要能依類型分類——內建、預設範例、"
             "自建，還有股市、美食這些主題——可以搜尋也可以篩選。」", 11, False, BODY)])

lbl2 = clone(s5, m["Text 7"][0], 0.90, 3.62, 0.40, 0.28, "Text S2 num")
write(lbl2, [("02", 14, True, GREEN)])
ttl2 = clone(s5, m["Text 8"][0], 1.30, 3.63, 6.00, 0.28, "Text S2 title")
write(ttl2, [("設定每天自動更新", 15, True, WHITE)])
lab2 = clone(s5, m["Text 0"][0], 1.30, 3.99, 5.00, 0.20, "Text S2 label")
mixed(lab2, [("DAILY ROUTINE", 8, False, FAINT, MONO), ("  ·  台北時間每天 09:00", 8, False, FAINT, CJK)])

steps = [
    ("每天 09:00 觸發", "台北時間，不必手動執行", AMBER),
    ("掃描 skill 目錄", "讀出分類、啟用狀態與說明", BLUE),
    ("與線上頁面比對", "直接比對 Artifact，不經過 GitHub", PURPLE),
    ("有變動才發佈", "沒變動就結束，不產生新版本", GREEN),
]
NW, GAP, TOP, NH = 2.70, 0.43, 4.40, 1.52
for i, (title, desc, color) in enumerate(steps):
    x = 0.62 + i * (NW + GAP)
    box = clone(s5, card, x, TOP, NW, NH, "Shape step %d" % (i + 1))
    dot = clone(s5, m["Text 7"][0], x + 0.24, TOP + 0.20, 0.40, 0.24, "Text step%d n" % (i + 1))
    write(dot, [("0%d" % (i + 1), 11, True, color, MONO)])
    tt = clone(s5, m["Text 8"][0], x + 0.24, TOP + 0.50, NW - 0.44, 0.28, "Text step%d t" % (i + 1))
    write(tt, [(title, 12.5, True, WHITE)])
    dd = clone(s5, m["Text 8"][0], x + 0.24, TOP + 0.82, NW - 0.44, 0.56, "Text step%d d" % (i + 1))
    write(dd, [(desc, 9.5, False, BODY)])
    if i < 3:
        ar = clone(s5, m["Text 0"][0], x + NW + 0.05, TOP + 0.62, 0.33, 0.24, "Text arrow%d" % (i + 1))
        write(ar, [("→", 13, False, FAINT, MONO)])

# Text 7 / Text 8 只是複製用的樣板，本體要移除，否則第五張會殘留第二張的文字
drop(m["Text 7"][0])
drop(m["Text 8"][0])

note = clone(s5, m["Text 0"][0], 0.62, 6.14, 12.09, 0.24, "Text note")
write(note, [("整條流程沒有人介入：掃描、比對、發佈都由排程完成，也不經過 GitHub。",
              10, False, FAINT)])

write(m["Text 16"][0], [("案例 · 製作流程", 9, False, FAINT)])

prs.save("out.pptx")
print("saved out.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")

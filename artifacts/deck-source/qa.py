# -*- coding: utf-8 -*-
"""從 pptx 讀出實際幾何：檢查出界、邊距、重疊，並輸出可在瀏覽器檢視的重建圖。"""
import base64, html, json, sys
from pptx import Presentation
from pptx.util import Emu

W, H = 13.333, 7.5
prs = Presentation("out.pptx")
issues = []
slides_html = []

def hexof(run):
    try:
        if run.font.color and run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return "EEF3FA"

for idx in (3, 4):
    s = prs.slides[idx]
    boxes = []
    parts = []
    for sh in s.shapes:
        if sh.left is None:
            continue
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        name = sh.name
        # 出界與邊距
        if l < -0.01 or t < -0.01 or l + w > W + 0.01 or t + h > H + 0.01:
            issues.append(f"slide{idx+1} 出界: {name} ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}")
        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        if txt and (l < 0.5 or l + w > W - 0.5):
            if not name.startswith("Text 2") :
                issues.append(f"slide{idx+1} 邊距不足: {name} 左{l:.2f} 右{W-(l+w):.2f}")
        if txt:
            boxes.append((name, l, t, w, h, txt))
        if sh.shape_type == 13:  # PICTURE
            img = base64.b64encode(sh.image.blob).decode()
            parts.append(f'<img style="left:{l/W*100:.3f}%;top:{t/H*100:.3f}%;'
                         f'width:{w/W*100:.3f}%;height:{h/H*100:.3f}%" src="data:image/png;base64,{img}">')
            continue
        if sh.has_text_frame and txt:
            runs = []
            for p in sh.text_frame.paragraphs:
                if not p.runs:
                    continue
                MONO_STACK = "Consolas,ui-monospace,monospace"
                CJK_STACK = "Noto Sans TC,sans-serif"
                chunks = []
                for r in p.runs:
                    fam = MONO_STACK if (r.font.name or "") == "Consolas" else CJK_STACK
                    sz = r.font.size.pt if r.font.size else 12
                    chunks.append(
                        '<span style="font-size:%.2fpt;font-weight:%d;color:#%s;font-family:%s">%s</span>'
                        % (sz, 700 if r.font.bold else 400, hexof(r), fam, html.escape(r.text)))
                inner = "".join(chunks)
                sa = p.space_after.pt if p.space_after else 0
                runs.append(f'<p style="margin:0 0 {sa:.1f}pt">{inner}</p>')
            parts.append(f'<div class="tb" style="left:{l/W*100:.3f}%;top:{t/H*100:.3f}%;'
                         f'width:{w/W*100:.3f}%;height:{h/H*100:.3f}%">{"".join(runs)}</div>')
        else:
            parts.append(f'<div class="sp" style="left:{l/W*100:.3f}%;top:{t/H*100:.3f}%;'
                         f'width:{w/W*100:.3f}%;height:{h/H*100:.3f}%"></div>')
    # 文字框兩兩重疊（面積重疊 > 0.05 平方吋才算）
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            ox = max(0, min(a[1]+a[3], b[1]+b[3]) - max(a[1], b[1]))
            oy = max(0, min(a[2]+a[4], b[2]+b[4]) - max(a[2], b[2]))
            if ox * oy > 0.05:
                issues.append(f"slide{idx+1} 文字框重疊 {ox*oy:.2f}in²: {a[0]}({a[5][:12]}) × {b[0]}({b[5][:12]})")
    slides_html.append(f'<section class="slide">{"".join(parts)}</section>')

open("qa.html", "w", encoding="utf-8").write("""<!doctype html><meta charset="utf-8">
<link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@400;500;700&display=swap">
<style>
body{margin:0;background:#333;font-family:'Noto Sans TC',sans-serif}
.slide{position:relative;width:1333px;height:750px;margin:16px auto;background:#0F131C;overflow:hidden}
.tb{position:absolute;overflow:visible;line-height:1.35}
.tb p{margin:0}
.sp{position:absolute;border:1px solid rgba(44,56,82,.9);border-radius:8px}
img{position:absolute;object-fit:fill}
</style>""" + "".join(slides_html))

print("=== 幾何檢查 ===")
print("\n".join(issues) if issues else "沒有出界、邊距不足或明顯重疊")
